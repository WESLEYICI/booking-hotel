from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x0B, 0x1D, 0x2F)
ACCENT = RGBColor(0x00, 0x9E, 0xE0)
LIGHT = RGBColor(0xE5, 0xF4, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
GRAY = RGBColor(0x60, 0x74, 0x8C)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x35)

def set_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color=ACCENT, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

# ─── SLIDE 1: COVER ───
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide1, DARK)

add_shape(slide1, Inches(0), Inches(2.5), Inches(13.333), Inches(0.06), ACCENT)

add_textbox(slide1, Inches(1), Inches(1.0), Inches(11), Inches(1.5),
            "MY HOTELS", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(2.7), Inches(11), Inches(1),
            "Modern Room Booking & Management System", font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(4.0), Inches(11), Inches(1),
            "Solusi Pemesanan Hotel Otomatis Terintegrasi", font_size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_shape(slide1, Inches(0), Inches(5.2), Inches(13.333), Inches(0.06), ACCENT)

add_textbox(slide1, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
            "Oleh: Tim Developer  •  Juni 2026", font_size=18, color=LIGHT, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: PROJECT CLIENT ───
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, WHITE)

add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_textbox(slide2, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
            "1. Project Client (Target Pengguna)", font_size=32, bold=True, color=WHITE)
add_textbox(slide2, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
            "Aplikasi website melayani dua sisi pengguna (Klien)", font_size=16, color=LIGHT)

# Left card - Hotel Management
card1 = add_shape(slide2, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0), RGBColor(0xF0, 0xF9, 0xFF))
add_textbox(slide2, Inches(1.2), Inches(1.8), Inches(4.5), Inches(0.5),
            "Pihak Hotel (Manajemen)", font_size=22, bold=True, color=DARK_TEXT)
add_bullet_textbox(slide2, Inches(1.2), Inches(2.5), Inches(4.8), Inches(3.5), [
    "Membutuhkan sistem digital untuk mendata kamar",
    "Memantau riwayat pesanan (Dashboard Analytics)",
    "Memvalidasi pembayaran secara otomatis",
    "Mengurangi human error operasional"
], font_size=16, color=DARK_TEXT)

# Right card - Guest
card2 = add_shape(slide2, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0), RGBColor(0xFD, 0xF4, 0xE4))
add_textbox(slide2, Inches(7.4), Inches(1.8), Inches(4.5), Inches(0.5),
            "Tamu (Guest)", font_size=22, bold=True, color=DARK_TEXT)
add_bullet_textbox(slide2, Inches(7.4), Inches(2.5), Inches(4.8), Inches(3.5), [
    "Melihat galeri kamar secara real-time",
    "Memfilter fasilitas kamar",
    "Melakukan pemesanan mandiri",
    "Membayar secara instan tanpa resepsionis"
], font_size=16, color=DARK_TEXT)

# ─── SLIDE 3: BPM / FLOWCHART ───
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3, WHITE)

add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_textbox(slide3, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
            "2. BPM / Flowchart (Alur Sistem)", font_size=32, bold=True, color=WHITE)
add_textbox(slide3, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
            "Alur pemesanan hotel dari awal hingga selesai", font_size=16, color=LIGHT)

# Flowchart - simplified as directional boxes with arrows
steps = [
    ("Mulai", ACCENT),
    ("Punya Akun?\nYa → Login\nTidak → Register", RGBColor(0x27, 0xAE, 0x60)),
    ("Pilih Kamar & Filter\nFasilitas", RGBColor(0x29, 0x80, 0xB9)),
    ("Isi Form Booking &\nMasukkan Voucher Diskon", RGBColor(0x8E, 0x44, 0xAD)),
    ("Proses Pembayaran\nvia Midtrans", RGBColor(0xE6, 0x7E, 0x22)),
    ("Status Pembayaran\n✅ Confirmed\n❌ Gagal/Batal", RGBColor(0xC0, 0x39, 0x2B)),
    ("Email Invoice\nOtomatis", RGBColor(0x27, 0xAE, 0x60)),
    ("Tamu Menginap &\nCheck-out", RGBColor(0x16, 0xA0, 0x85)),
    ("Beri Ulasan/Review\nBintang", GOLD),
    ("Selesai", ACCENT),
]

y_start = 1.6
for i, (text, color) in enumerate(steps):
    col = i // 5
    row = i % 5
    x = Inches(0.8 + col * 6.2)
    y = Inches(y_start + row * 1.05)
    s = add_shape(slide3, x, y, Inches(2.6), Inches(0.85), color)
    add_textbox(slide3, x + Inches(0.05), y + Inches(0.05), Inches(2.5), Inches(0.75),
                text, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE 4: TECH STACK ───
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide4, WHITE)

add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_textbox(slide4, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
            "3. Teknologi yang Digunakan (Tech Stack)", font_size=32, bold=True, color=WHITE)
add_textbox(slide4, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
            "Full-Stack JavaScript (MERN/PERN Stack variation)", font_size=16, color=LIGHT)

sections = [
    ("Frontend (Antarmuka)", RGBColor(0x29, 0x80, 0xB9), [
        "React.js — Framework utama UI reaktif",
        "Tailwind CSS & Flowbite — Desain modern & responsif",
    ]),
    ("Backend (Server & Database)", RGBColor(0x8E, 0x44, 0xAD), [
        "Node.js & Express.js — RESTful API server",
        "MySQL — Database relasional solid",
    ]),
    ("Layanan Pihak Ketiga (3rd Party / API)", RGBColor(0xE6, 0x7E, 0x22), [
        "Midtrans Payment Gateway — QRIS, Transfer, E-Wallet",
        "NodeMailer (Google SMTP) — Email invoice otomatis",
    ]),
]

for i, (title, accent, items) in enumerate(sections):
    x = Inches(0.6 + i * 4.2)
    card = add_shape(slide4, x, Inches(1.6), Inches(3.8), Inches(5.2), RGBColor(0xF8, 0xF9, 0xFA))
    hdr = add_shape(slide4, x, Inches(1.6), Inches(3.8), Inches(0.7), accent)
    add_textbox(slide4, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.5),
                title, font_size=16, bold=True, color=WHITE)
    add_bullet_textbox(slide4, x + Inches(0.2), Inches(2.6), Inches(3.4), Inches(3.5),
                        items, font_size=14, color=DARK_TEXT, spacing=Pt(12))

# ─── SLIDE 5: GITHUB ───
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide5, WHITE)

add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_textbox(slide5, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
            "4. GitHub (Repository)", font_size=32, bold=True, color=WHITE)
add_textbox(slide5, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
            "Version Control System Git", font_size=16, color=LIGHT)

card = add_shape(slide5, Inches(2), Inches(2.0), Inches(9), Inches(3.5), RGBColor(0xF0, 0xF9, 0xFF))
add_textbox(slide5, Inches(2.5), Inches(2.3), Inches(8), Inches(0.5),
            "Link Repository", font_size=22, bold=True, color=DARK_TEXT)

link_shape = add_shape(slide5, Inches(2.5), Inches(3.0), Inches(8), Inches(0.7), ACCENT)
add_textbox(slide5, Inches(2.7), Inches(3.05), Inches(7.5), Inches(0.6),
            "https://github.com/WESLEYICI/booking-hotel.git", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_bullet_textbox(slide5, Inches(2.5), Inches(4.0), Inches(8), Inches(1.5), [
    "Cabang Utama: main",
    "Keamanan: Dilengkapi Secret Scanning untuk melindungi kredensial Payment Gateway"
], font_size=14, color=DARK_TEXT)

# ─── SLIDE 6: ROLE USER ───
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide6, WHITE)

add_shape(slide6, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_textbox(slide6, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
            "5. Role User (Peran Pengguna)", font_size=32, bold=True, color=WHITE)
add_textbox(slide6, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
            "2 jenis Hak Akses dalam aplikasi", font_size=16, color=LIGHT)

# Guest card
card_g = add_shape(slide6, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3), RGBColor(0xEB, 0xF5, 0xFB))
add_shape(slide6, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.65), RGBColor(0x29, 0x80, 0xB9))
add_textbox(slide6, Inches(1.0), Inches(1.7), Inches(5.0), Inches(0.5),
            "Guest (Tamu)", font_size=22, bold=True, color=WHITE)
add_bullet_textbox(slide6, Inches(1.0), Inches(2.5), Inches(5.0), Inches(4.0), [
    "Mencari & memfilter daftar kamar tersedia",
    "Booking & memasukkan kode Promo/Voucher",
    "Bayar via Midtrans (VA, E-Wallet)",
    "Melihat riwayat (My Bookings) + Email invoice",
    "Memberi ulasan & rating kamar",
    "Live Chat Bantuan via WhatsApp"
], font_size=14, color=DARK_TEXT)

# Admin card
card_a = add_shape(slide6, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.3), RGBColor(0xFD, 0xF2, 0xE9))
add_shape(slide6, Inches(6.9), Inches(1.6), Inches(5.8), Inches(0.65), RGBColor(0xE6, 0x7E, 0x22))
add_textbox(slide6, Inches(7.3), Inches(1.7), Inches(5.0), Inches(0.5),
            "Admin (Manajemen Hotel)", font_size=22, bold=True, color=WHITE)
add_bullet_textbox(slide6, Inches(7.3), Inches(2.5), Inches(5.0), Inches(4.0), [
    "Dashboard Analytics (pesanan, pendapatan)",
    "CRUD Rooms — Tambah, edit, nonaktifkan kamar",
    "Manage Bookings — Riwayat + update status",
    "Voucher Management — Buat kode diskon",
    "Review Monitoring — Baca ulasan pelanggan"
], font_size=14, color=DARK_TEXT)

# ─── SLIDE 7: DEMO VIDEO ───
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide7, DARK)

add_shape(slide7, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_textbox(slide7, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
            "6. Demo Aplikasi", font_size=32, bold=True, color=WHITE)
add_textbox(slide7, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
            "Video Screen Recording (1-2 menit)", font_size=16, color=LIGHT)

# Placeholder for video
vid_placeholder = add_shape(slide7, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.5), RGBColor(0x34, 0x40, 0x50))
add_textbox(slide7, Inches(2), Inches(2.5), Inches(9.5), Inches(1),
            "▶  INSERT VIDEO HERE", font_size=40, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide7, Inches(2), Inches(3.5), Inches(9.5), Inches(0.5),
            "Screen Recording Demo Aplikasi", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Scenario list
add_shape(slide7, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.04), ACCENT)
add_textbox(slide7, Inches(1.5), Inches(5.8), Inches(10.5), Inches(0.4),
            "Skenario Demo:", font_size=16, bold=True, color=ACCENT)
add_bullet_textbox(slide7, Inches(1.5), Inches(6.15), Inches(10.5), Inches(1.2), [
    "Halaman depan (Jumbotron) → Rooms & Suites → Filter Fasilitas/Harga",
    "Detail kamar + Ulasan Tamu → Booking + Voucher → Bayar via Midtrans",
    "Email Invoice otomatis → Login Admin → Dashboard Data Bookings"
], font_size=13, color=LIGHT)

# ─── LAST SLIDE: THANK YOU ───
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide8, DARK)

add_shape(slide8, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), GOLD)
add_textbox(slide8, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "Terima Kasih", font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide8, Inches(1), Inches(3.4), Inches(11), Inches(0.8),
            "Any Questions?", font_size=28, color=ACCENT, align=PP_ALIGN.CENTER)
add_shape(slide8, Inches(0), Inches(4.6), Inches(13.333), Inches(0.06), GOLD)

add_textbox(slide8, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
            "github.com/WESLEYICI/booking-hotel.git", font_size=16, color=LIGHT, align=PP_ALIGN.CENTER)

output_path = r"C:\laragon\www\booking-hotel\My_Hotels_Presentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
